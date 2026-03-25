#!/usr/bin/env python3
"""Generate launch-manifest-overview.pptx from template.

Usage: python3 docs/presentations/gen_manifest_slides.py

Reads:  docs/presentations/launch-manifest-overview.pptx (2-slide template)
Writes: docs/presentations/launch-manifest-overview.pptx (full deck)
"""

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import textwrap

TEMPLATE = "docs/presentations/launch-manifest-overview.pptx"
OUTPUT = "docs/presentations/launch-manifest-overview.pptx"

# Layout indices — "Simple Light" 2nd group (matches template slides)
L_TITLE = 11        # Simple Light / TITLE
L_BODY = 12         # Simple Light / TITLE_AND_BODY

# Content slide title font size
TITLE_SIZE = Pt(24)

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x1A, 0x73, 0xE8)  # Blue
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)  # Dark code background
CODE_FG = RGBColor(0xCD, 0xD6, 0xF4)  # Light code text
CODE_KEY = RGBColor(0x89, 0xB4, 0xFA)  # Blue for YAML keys
CODE_STR = RGBColor(0xA6, 0xE3, 0xA1)  # Green for strings/values
CODE_CMT = RGBColor(0x6C, 0x70, 0x86)  # Gray for comments
GREEN = RGBColor(0x34, 0xA8, 0x53)
RED = RGBColor(0xEA, 0x43, 0x35)
ORANGE = RGBColor(0xFB, 0xBC, 0x04)

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def remove_all_slides(prs):
    """Remove all slides from presentation."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def remove_bullets(para):
    """Remove bullet from a paragraph."""
    pPr = para._p.find(f"{{{NS_A}}}pPr")
    if pPr is None:
        pPr = etree.SubElement(para._p, f"{{{NS_A}}}pPr")
        para._p.insert(0, pPr)
    # Remove existing bullet elements
    for tag in ["buChar", "buAutoNum", "buFont", "buSzPct", "buSzPts"]:
        for el in pPr.findall(f"{{{NS_A}}}{tag}"):
            pPr.remove(el)
    # Add buNone
    existing = pPr.find(f"{{{NS_A}}}buNone")
    if existing is None:
        etree.SubElement(pPr, f"{{{NS_A}}}buNone")


def add_run(para, text, size=None, bold=None, italic=None, color=None, font_name=None):
    """Add a styled run to a paragraph."""
    run = para.add_run()
    run.text = text
    if size:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    return run


def add_code_block(slide, left, top, width, height, lines):
    """Add a colored code block as a text box with dark background.

    lines is a list of (text, color) tuples or strings (default color).
    Each item is one line. Items can also be lists of (text, color) segments.
    """
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    # Dark background
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = CODE_BG

    # Rounded corners via manual XML
    spPr = shape._element.find(f"{{{NS_A}}}..//{{{NS_A}}}../")

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.space_before = Pt(0)
        para.space_after = Pt(0)
        para.line_spacing = Pt(14)
        remove_bullets(para)

        if isinstance(line, str):
            add_run(para, line, size=Pt(10), color=CODE_FG, font_name="Consolas")
        elif isinstance(line, list):
            for text, color in line:
                add_run(para, text, size=Pt(10), color=color, font_name="Consolas")
        elif isinstance(line, tuple):
            text, color = line
            add_run(para, text, size=Pt(10), color=color, font_name="Consolas")


def yaml_line(key, value="", comment=""):
    """Build a colored YAML line as segments."""
    parts = []
    if key:
        parts.append((key + ":", CODE_KEY))
        if value:
            parts.append((" " + value, CODE_STR))
    elif value:
        parts.append((value, CODE_FG))
    if comment:
        parts.append(("  " + comment, CODE_CMT))
    return parts


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    slide.placeholders[0].text = title
    # Template subtitle placeholder is unreliable — use manual textbox
    if subtitle:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(8.0), Inches(0.8))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        remove_bullets(p)
        add_run(p, subtitle, size=Pt(16), color=GRAY)
    return slide


def add_body_slide(prs, title, bullets=None, no_bullets=False):
    """Add a slide with title and optional bullet list.

    bullets: list of strings or (text, level) tuples.
    Returns (slide, body_placeholder) for further customization.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[L_BODY])
    # Set title with explicit size
    tf = slide.placeholders[0].text_frame
    tf.paragraphs[0].text = ""
    add_run(tf.paragraphs[0], title, size=TITLE_SIZE, bold=True, color=NEAR_BLACK)

    body = slide.placeholders[1]
    if bullets:
        tf = body.text_frame
        for i, item in enumerate(bullets):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            if isinstance(item, tuple):
                text, level = item
                para.level = level
            else:
                text = item
                para.level = 0

            if no_bullets:
                remove_bullets(para)

            para.text = text

    return slide, body


def add_two_col_slide(prs, title, left_items, right_items):
    """Two-column slide using body layout + manual text boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_BODY])
    tf = slide.placeholders[0].text_frame
    tf.paragraphs[0].text = ""
    add_run(tf.paragraphs[0], title, size=TITLE_SIZE, bold=True, color=NEAR_BLACK)
    # Clear body placeholder
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])

    col_top = Inches(1.1)
    col_height = Inches(3.8)

    for col_left, col_width, items in [
        (Inches(0.5), Inches(4.3), left_items),
        (Inches(5.0), Inches(4.5), right_items),
    ]:
        box = slide.shapes.add_textbox(col_left, col_top, col_width, col_height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            if isinstance(item, tuple):
                text, level = item
            else:
                text = item
                level = 0
            if level == 0:
                add_run(para, text, size=Pt(16), bold=True, color=NEAR_BLACK)
            else:
                add_run(para, "  " + text, size=Pt(14), color=GRAY)
            para.space_before = Pt(2)
            para.space_after = Pt(2)
            remove_bullets(para)

    return slide


# ─── Body area coordinates (from template inspection) ───
BODY_LEFT = Inches(0.53)
BODY_TOP = Inches(1.18)
BODY_WIDTH = Inches(9.0)
BODY_HEIGHT = Inches(4.1)

# Code block defaults
CODE_LEFT = Inches(0.6)
CODE_WIDTH = Inches(8.7)


def main():
    prs = Presentation(TEMPLATE)
    remove_all_slides(prs)

    # ═══════════════════════════════════════════════════════
    # Slide 1: Title
    # ═══════════════════════════════════════════════════════
    add_title_slide(prs, "Launch Manifest",
                    "Declaring and Auditing the ROS 2 Communication Graph\nplay_launch — Phase 31")

    # ═══════════════════════════════════════════════════════
    # Slide 2: Problem
    # ═══════════════════════════════════════════════════════
    add_body_slide(prs, "Problem", [
        "Launch files declare nodes, not topics",
        "Topic creation happens in source code — invisible until runtime",
        "Autoware: ~110 nodes create ~500 topics",
        "A code change adds/removes a topic → goes unnoticed",
        'No way to answer: "what SHOULD the graph look like?"',
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 3: Manifest in 30 Seconds
    # ═══════════════════════════════════════════════════════
    slide, _ = add_body_slide(prs, "Manifest in 30 Seconds")
    # Clear body placeholder
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])

    add_code_block(slide, CODE_LEFT, Inches(1.0), CODE_WIDTH, Inches(3.5), [
        [("version", CODE_KEY), (": 1", CODE_STR)],
        "",
        [("nodes", CODE_KEY), (":", CODE_KEY)],
        [("  talker", CODE_KEY), (":", CODE_KEY)],
        [("    pub", CODE_KEY), (": [chatter]", CODE_STR)],
        [("  listener", CODE_KEY), (":", CODE_KEY)],
        [("    sub", CODE_KEY), (": [chatter]", CODE_STR)],
        "",
        [("topics", CODE_KEY), (":", CODE_KEY)],
        [("  chatter", CODE_KEY), (":", CODE_KEY)],
        [("    type", CODE_KEY), (": std_msgs/msg/String", CODE_STR)],
        [("    pub", CODE_KEY), (": [talker/chatter]", CODE_STR)],
        [("    sub", CODE_KEY), (": [listener/chatter]", CODE_STR)],
        "",
        [("exports", CODE_KEY), (":", CODE_KEY)],
        [("  output", CODE_KEY), (": [talker/chatter]", CODE_STR)],
    ])

    # Caption below code
    cap = slide.shapes.add_textbox(CODE_LEFT, Inches(4.55), CODE_WIDTH, Inches(0.4))
    p = cap.text_frame.paragraphs[0]
    remove_bullets(p)
    add_run(p, "One manifest per launch file. Topics wire endpoints. Names are relative.", size=Pt(14), color=GRAY)

    # ═══════════════════════════════════════════════════════
    # Slide 4: Manifest Structure
    # ═══════════════════════════════════════════════════════
    add_body_slide(prs, "Manifest Structure", [
        "nodes:        endpoint declarations (pre-remap names)",
        "topics:        wiring + type + QoS + channel contract",
        "includes:     child scopes (name = namespace, manifest path)",
        "imports:      subscriber dependencies (what we need from outside)",
        "exports:      publisher outputs (what we provide to parent)",
        "io:              scope-level timing contract",
    ], no_bullets=True)

    # ═══════════════════════════════════════════════════════
    # Slide 5: Wiring
    # ═══════════════════════════════════════════════════════
    slide, _ = add_body_slide(prs, "Wiring: Topics Connect Endpoints")
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])

    add_code_block(slide, CODE_LEFT, Inches(1.0), Inches(5.2), Inches(3.0), [
        [("nodes", CODE_KEY), (":", CODE_KEY)],
        [("  cropbox_filter", CODE_KEY), (":", CODE_KEY)],
        [("    pub", CODE_KEY), (": [output]", CODE_STR),
         ("       # pre-remap name", CODE_CMT)],
        [("    sub", CODE_KEY), (": [input]", CODE_STR)],
        "",
        [("topics", CODE_KEY), (":", CODE_KEY)],
        [("  cropped", CODE_KEY), (":", CODE_KEY),
         ("             # → <ns>/cropped", CODE_CMT)],
        [("    type", CODE_KEY), (": PointCloud2", CODE_STR)],
        [("    pub", CODE_KEY), (": [cropbox_filter/output]", CODE_STR)],
        [("    sub", CODE_KEY), (": [ground_filter/input]", CODE_STR)],
        [("    contract", CODE_KEY), (":", CODE_KEY)],
        [("      rate_hz", CODE_KEY), (": 10", CODE_STR)],
    ])

    # Side notes
    notes = slide.shapes.add_textbox(Inches(6.2), Inches(1.2), Inches(3.3), Inches(2.8))
    tf = notes.text_frame
    tf.word_wrap = True
    items = [
        ("Endpoint names", " = node's internal topic names (before <remap>)"),
        ("Topic", " gives the real name and wires pub ↔ sub"),
        ("contract:", " = channel-level quality bound (rate, deadline, drops)"),
    ]
    for i, (bold_text, normal_text) in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(8)
        add_run(para, bold_text, size=Pt(13), bold=True, color=ACCENT)
        add_run(para, normal_text, size=Pt(13), color=NEAR_BLACK)

    # ═══════════════════════════════════════════════════════
    # Slide 6: Composition
    # ═══════════════════════════════════════════════════════
    slide, _ = add_body_slide(prs, "Composition: Includes and Scope Hierarchy")
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])

    add_code_block(slide, CODE_LEFT, Inches(1.0), CODE_WIDTH, Inches(3.2), [
        [("includes", CODE_KEY), (":", CODE_KEY)],
        [("  lidar", CODE_KEY), (":", CODE_KEY),
         ("                        # name = ROS namespace", CODE_CMT)],
        [("    manifest", CODE_KEY), (": tier4_.../lidar_perception.launch.yaml", CODE_STR)],
        [("  camera", CODE_KEY), (":", CODE_KEY)],
        [("    manifest", CODE_KEY), (": tier4_.../camera_perception.launch.yaml", CODE_STR)],
        "",
        [("topics", CODE_KEY), (":", CODE_KEY)],
        [("  lidar_objects", CODE_KEY), (":", CODE_KEY)],
        [("    pub", CODE_KEY), (": [lidar/output]", CODE_STR),
         ("          # child's export", CODE_CMT)],
        [("    sub", CODE_KEY), (": [fusion_node/lidar_in]", CODE_STR)],
        "",
        [("imports", CODE_KEY), (":", CODE_KEY)],
        [("  input", CODE_KEY), (": [lidar/input, camera/input]", CODE_STR)],
        [("exports", CODE_KEY), (":", CODE_KEY)],
        [("  output", CODE_KEY), (": [tracker/tracked]", CODE_STR)],
    ])

    cap = slide.shapes.add_textbox(CODE_LEFT, Inches(4.3), CODE_WIDTH, Inches(0.4))
    p = cap.text_frame.paragraphs[0]
    remove_bullets(p)
    add_run(p, "Parent wires children explicitly via topics. No implicit accumulation.", size=Pt(14), bold=True, color=NEAR_BLACK)

    # ═══════════════════════════════════════════════════════
    # Slide 7: I/O Patterns
    # ═══════════════════════════════════════════════════════
    add_two_col_slide(prs, "I/O Patterns in Autoware", [
        "Pipe (1→1)",
        ("on_arrival — cropbox, centerpoint", 1),
        "",
        "Fusion (N→1)",
        ("all_ready + timestamp — concat_pointclouds", 1),
        "",
        "Source (0→N)",
        ("periodic — lidar_driver (10 Hz)", 1),
    ], [
        "Timer-driven",
        ("periodic + state inputs — tracker (10 Hz)", 1),
        "",
        "Feedback cycle",
        ("on_arrival + state — NDT ↔ EKF", 1),
        "",
        "",
        "All from the same primitives —",
        "no enumerated pattern types.",
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 8: Node I/O Contract
    # ═══════════════════════════════════════════════════════
    slide, _ = add_body_slide(prs, "Node I/O Contract")
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])

    add_code_block(slide, CODE_LEFT, Inches(1.0), Inches(4.8), Inches(3.8), [
        ("# Simple pipe", CODE_CMT),
        [("centerpoint", CODE_KEY), (":", CODE_KEY)],
        [("  io", CODE_KEY), (": { ", CODE_FG), ("latency_ms", CODE_KEY),
         (": 30 }", CODE_STR)],
        "",
        ("# Fusion — all inputs synchronized", CODE_CMT),
        [("fusion_node", CODE_KEY), (":", CODE_KEY)],
        [("  io", CODE_KEY), (":", CODE_KEY)],
        [("    trigger", CODE_KEY), (": all_ready", CODE_STR)],
        [("    correlation", CODE_KEY), (": timestamp", CODE_STR)],
        [("    tolerance_ms", CODE_KEY), (": 50", CODE_STR)],
        [("    latency_ms", CODE_KEY), (": 20", CODE_STR)],
        "",
        ("# Timer-driven with state inputs", CODE_CMT),
        [("tracker", CODE_KEY), (":", CODE_KEY)],
        [("  io", CODE_KEY), (":", CODE_KEY)],
        [("    state", CODE_KEY), (": [fused]", CODE_STR)],
        [("    trigger", CODE_KEY), (": { ", CODE_FG), ("periodic", CODE_KEY),
         (": 10 }", CODE_STR)],
        [("    freshness_ms", CODE_KEY), (": 200", CODE_STR)],
    ])

    # Side note
    notes = slide.shapes.add_textbox(Inches(5.8), Inches(1.2), Inches(3.7), Inches(2.0))
    tf = notes.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "io:", size=Pt(14), bold=True, color=ACCENT, font_name="Consolas")
    add_run(p, " only adds annotations to ", size=Pt(14), color=NEAR_BLACK)
    add_run(p, "sub:/pub:", size=Pt(14), bold=True, color=ACCENT, font_name="Consolas")
    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    add_run(p2, "No redundancy — timing and roles are separate from topology.", size=Pt(14), color=NEAR_BLACK)

    # ═══════════════════════════════════════════════════════
    # Slide 9: Contract Primitives
    # ═══════════════════════════════════════════════════════
    add_two_col_slide(prs, "Contract Primitives", [
        "Input roles (per sub: endpoint)",
        ("role: trigger (default) — fires computation", 1),
        ("role: state — read-latest when triggered", 1),
        ("consume: N — messages per firing", 1),
        ("min_count: N — readiness precondition", 1),
        "",
        "Trigger",
        ("on_arrival — any trigger-input (default)", 1),
        ("all_ready — barrier, all must arrive", 1),
        ("periodic: N — timer at N Hz", 1),
    ], [
        "Correlation (2+ trigger inputs)",
        ("timestamp — match by header.stamp", 1),
        ("latest — most recent of each", 1),
        "",
        "Timing bounds",
        ("latency_ms — trigger → output", 1),
        ("freshness_ms — max input data age", 1),
        "",
        "Drop tolerance",
        ("max_consecutive_drops", 1),
        ("max_drop_ratio over window", 1),
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 10: Topic Channel Contract
    # ═══════════════════════════════════════════════════════
    slide, _ = add_body_slide(prs, "Topic Channel Contract")
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])

    add_code_block(slide, CODE_LEFT, Inches(1.0), Inches(5.0), Inches(2.2), [
        [("topics", CODE_KEY), (":", CODE_KEY)],
        [("  detected_objects", CODE_KEY), (":", CODE_KEY)],
        [("    type", CODE_KEY), (": DetectedObjects", CODE_STR)],
        [("    pub", CODE_KEY), (": [centerpoint/objects]", CODE_STR)],
        [("    sub", CODE_KEY), (": [tracker/input]", CODE_STR)],
        [("    contract", CODE_KEY), (":", CODE_KEY)],
        [("      rate_hz", CODE_KEY), (": 10", CODE_STR)],
        [("      deadline_ms", CODE_KEY), (": 150", CODE_STR)],
        [("      max_drop_ratio", CODE_KEY), (": 0.05", CODE_STR)],
    ])

    notes = slide.shapes.add_textbox(CODE_LEFT, Inches(3.4), CODE_WIDTH, Inches(1.2))
    tf = notes.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_run(p, "Two separate concerns:", size=Pt(16), bold=True, color=NEAR_BLACK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    add_run(p2, "Node io:  ", size=Pt(14), bold=True, color=ACCENT, font_name="Consolas")
    add_run(p2, "computation — how fast input → output", size=Pt(14), color=NEAR_BLACK)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(4)
    add_run(p3, "Topic contract:  ", size=Pt(14), bold=True, color=ACCENT, font_name="Consolas")
    add_run(p3, "channel — what the link promises", size=Pt(14), color=NEAR_BLACK)

    # ═══════════════════════════════════════════════════════
    # Slide 11: Full Pipeline (code example)
    # ═══════════════════════════════════════════════════════
    slide, _ = add_body_slide(prs, "Full Pipeline: Top-Level Manifest")
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])

    add_code_block(slide, CODE_LEFT, Inches(1.0), Inches(5.5), Inches(3.5), [
        [("includes", CODE_KEY), (":", CODE_KEY)],
        [("  sensing", CODE_KEY), (":", CODE_KEY)],
        [("    manifest", CODE_KEY), (": tier4_.../sensing.launch.yaml", CODE_STR)],
        [("  perception", CODE_KEY), (":", CODE_KEY)],
        [("    manifest", CODE_KEY), (": tier4_.../perception.launch.yaml", CODE_STR)],
        [("  planning", CODE_KEY), (":", CODE_KEY)],
        [("    manifest", CODE_KEY), (": tier4_.../planning.launch.yaml", CODE_STR)],
        "",
        [("topics", CODE_KEY), (":", CODE_KEY)],
        [("  pointcloud", CODE_KEY), (":", CODE_KEY)],
        [("    pub", CODE_KEY), (": [sensing/pointcloud]", CODE_STR)],
        [("    sub", CODE_KEY), (": [perception/input]", CODE_STR)],
        [("  tracked_objects", CODE_KEY), (":", CODE_KEY)],
        [("    pub", CODE_KEY), (": [perception/output]", CODE_STR)],
        [("    sub", CODE_KEY), (": [planning/input]", CODE_STR)],
    ])

    # Side: latency summary
    notes = slide.shapes.add_textbox(Inches(6.4), Inches(1.0), Inches(3.2), Inches(3.0))
    tf = notes.text_frame
    tf.word_wrap = True
    items = [
        ("Scope budgets:", True, NEAR_BLACK),
        ("  sensing: source, 10 Hz", False, GRAY),
        ("  perception: 85 ms", False, GRAY),
        ("    lidar: 50 ms", False, GRAY),
        ("    camera: 30 ms", False, GRAY),
        ("  planning: 100 ms", False, GRAY),
        ("", False, GRAY),
        ("Critical path:", True, NEAR_BLACK),
        ("max(50,30) + 20 + 100", False, GRAY),
    ]
    for i, (text, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        remove_bullets(p)
        add_run(p, text, size=Pt(13), bold=bold, color=color)
    p = tf.add_paragraph()
    remove_bullets(p)
    add_run(p, "= 170 ms", size=Pt(16), bold=True, color=RED)

    # ═══════════════════════════════════════════════════════
    # Slide 12: Working with Launch Files
    # ═══════════════════════════════════════════════════════
    add_body_slide(prs, "Working with Launch Files", [
        "Manifest loaded alongside <include> — parser resolves namespace",
        "Endpoint names match <remap from=\"...\"> in launch XML",
        ("Warn if remap is missing", 1),
        "Same manifest reused across different namespace contexts",
        "Partial coverage: missing manifest = no auditing for that scope",
        "Full ROS 2 compatibility — manifests are sidecar files, not replacements",
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 13: Contract Violation Detection (diagram)
    # ═══════════════════════════════════════════════════════
    slide, _ = add_body_slide(prs, "Contract Violation Detection via RCL Interception")
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])
    slide.shapes.add_picture("docs/presentations/img/rcl-interception.png",
                             Inches(0.3), Inches(1.0), Inches(9.3))

    # ═══════════════════════════════════════════════════════
    # Slide 14: What the Monitor Checks
    # ═══════════════════════════════════════════════════════
    slide, _ = add_body_slide(prs, "What the Monitor Checks")
    slide.placeholders[1].text_frame.paragraphs[0].text = ""
    remove_bullets(slide.placeholders[1].text_frame.paragraphs[0])

    # Left: check list as textbox
    checks = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(4.5), Inches(3.5))
    tf = checks.text_frame
    tf.word_wrap = True
    items = [
        ("Node io: checks", True),
        ("  latency_ms: pub_time - take_time ≤ bound", False),
        ("  freshness_ms: pub_time - stamp ≤ bound", False),
        ("", False),
        ("Topic contract: checks", True),
        ("  rate_hz: interval ≈ 1/rate", False),
        ("  deadline_ms: interval ≤ deadline", False),
        ("  max_drop_ratio: drops/total ≤ ratio", False),
    ]
    for i, (text, is_bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        remove_bullets(p)
        add_run(p, text, size=Pt(14), bold=is_bold,
                color=NEAR_BLACK if is_bold else GRAY)

    # Right: 4-way diagnosis table
    rows, cols = 5, 3
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(5.0), Inches(1.0), Inches(4.5), Inches(2.5))
    tbl = tbl_shape.table

    # Header
    header_data = ["Assumption", "Guarantee", "Diagnosis"]
    for c, text in enumerate(header_data):
        cell = tbl.cell(0, c)
        cell.text = text
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x58, 0x58, 0x58)

    # Data rows
    data = [
        ("OK", "OK", "Nominal", GREEN, GREEN, GREEN),
        ("OK", "FAIL", "Node bug", GREEN, RED, RED),
        ("FAIL", "OK", "Node robust", RED, GREEN, ACCENT),
        ("FAIL", "FAIL", "Upstream problem", RED, RED, ORANGE),
    ]
    for r, (a, g, diag, a_color, g_color, d_color) in enumerate(data, 1):
        for c, (text, color) in enumerate([(a, a_color), (g, g_color), (diag, d_color)]):
            cell = tbl.cell(r, c)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = color
            # Light background for OK/FAIL cells
            if c < 2:
                cell.fill.solid()
                if color == GREEN:
                    cell.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xEA)
                else:
                    cell.fill.fore_color.rgb = RGBColor(0xFC, 0xE8, 0xE6)

    # ═══════════════════════════════════════════════════════
    # Slide 15: Scope Budget Verification
    # ═══════════════════════════════════════════════════════
    add_body_slide(prs, "Scope Budget Verification", [
        "Critical path = longest path in DAG → O(V+E)",
        "",
        "scope.io.latency_ms ≥ longest_path(imports → exports)",
        "",
        "Composition rules:",
        ("Series: sum of latencies", 1),
        ("Parallel: max of branches", 1),
        ("Periodic: period + jitter", 1),
        "",
        "Checked statically at parse time — before any node starts",
        "Algorithm: topological sort + dynamic programming (petgraph)",
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 16: Contract Theory
    # ═══════════════════════════════════════════════════════
    add_body_slide(prs, "Contract Theory Overview", [
        "Assume-Guarantee: C = (A, G)",
        ("A = input constraints (rate, availability)", 1),
        ("G = output promises (latency, rate) — valid only when A holds", 1),
        "",
        "Refinement: C' refines C ⟺ weaker A' + stronger G'",
        ("Faster node = valid upgrade", 1),
        ("More tolerant node = valid upgrade", 1),
        "",
        "QoS compatibility as base contract layer",
        ("reliability ≥, durability ≥, deadline ≤", 1),
        "",
        "Empirical derivation: capture → observe → derive contracts",
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 17: Verification Tooling
    # ═══════════════════════════════════════════════════════
    add_body_slide(prs, "Verification Tooling Path", [
        "Tier 1 — Now (covers ~95% of checks)",
        ("petgraph + plain Rust", 1),
        ("Critical path, scope budgets, QoS, types, wiring", 1),
        ("Zero new dependencies", 1),
        "",
        "Tier 2 — With audit feature",
        ("Hand-rolled runtime monitors + interception events", 1),
        ("Rate, deadline, latency, drops", 1),
        "",
        "Tier 3 — If needed",
        ("z3: constraint satisfiability for large manifests", 1),
        ("good_lp: budget optimization across pipelines", 1),
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 18: Workflow
    # ═══════════════════════════════════════════════════════
    add_body_slide(prs, "Workflow", [
        "1. CAPTURE",
        ("play_launch launch <pkg> <file> --save-manifest-dir manifests/", 1),
        ("Run → stabilize → snapshot → per-launch-file manifests", 1),
        "",
        "2. REFINE",
        ("Edit manifests: add io: contracts, tighten bounds", 1),
        "",
        "3. AUDIT",
        ("play_launch launch <pkg> <file> --manifest-dir manifests/", 1),
        ("Parse-time: scope budgets, QoS, wiring checks", 1),
        ("Runtime: latency, rate, drops via RCL interception", 1),
        "",
        "Gradual adoption: start with topology, add timing later",
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 19: Status
    # ═══════════════════════════════════════════════════════
    add_two_col_slide(prs, "Status and Next Steps", [
        "Done",
        ("Manifest format design (v1)", 1),
        ("I/O contract primitives", 1),
        ("RCL interception (Phase 29)", 1),
        ("Scope table in record.json (Phase 30)", 1),
        ("Contract theory formalization", 1),
        ("Prior art survey", 1),
    ], [
        "Next",
        ("Manifest parser (Rust crate)", 1),
        ("Parse-time consistency checks", 1),
        ("Capture mode (graph → manifests)", 1),
        ("Runtime audit (interception → monitor)", 1),
        ("Web UI manifest diff view", 1),
    ])

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
